import pptx
from pptx.dml.color import RGBColor
import os

def change_presentation_background(input_path, output_path):
    print(f"Loading presentation: {input_path}")
    prs = pptx.Presentation(input_path)
    
    # A professional, sleek dark blue/slate color (commonly used in modern tech presentations)
    # Hex: #0B192C -> RGB: 11, 25, 44
    bg_color = RGBColor(11, 25, 44)
    
    # Text colors need to be adjusted if the background becomes very dark.
    # We will try to find text and make it white/light gray if it's too dark.
    light_text_color = RGBColor(240, 240, 240)
    
    print("Applying new background to all slides...")
    for i, slide in enumerate(prs.slides):
        # 1. Change background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        # 2. Adjust text colors to stay visible on dark background
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
                
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # If the text doesn't have an explicit color, or if it's set to something very dark,
                    # we should force it to be light so it's readable on the new dark background.
                    # python-pptx can be tricky with theme colors, but we'll apply a blanket white/off-white
                    # to make it pop against the new dark background.
                    run.font.color.rgb = light_text_color

    print(f"Saving modified presentation to: {output_path}")
    prs.save(output_path)
    print("Done!")

if __name__ == "__main__":
    base_dir = r"c:\Users\chahi\Downloads\intern-p-main"
    input_file = os.path.join(base_dir, "Intelligence (1).pptx")
    output_file = os.path.join(base_dir, "Intelligence_Updated_Dark.pptx")
    
    change_presentation_background(input_file, output_file)
